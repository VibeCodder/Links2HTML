import sys
import re
import html
from PyQt6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout,
                             QHBoxLayout, QTextEdit, QPlainTextEdit, QPushButton,
                             QFileDialog, QLabel, QMessageBox)
from bs4 import BeautifulSoup
import mammoth
from pptx import Presentation
from pptx.oxml.ns import qn
 
# Regular expression catching raw links (http/https) in plain text
URL_REGEX = re.compile(r'(https?://[^\s<()\"\']+)')

# Regular expression catching an *already existing* literal HTML link, e.g.
# when the user pastes raw HTML source such as <a href="...">text</a> into
# the preview area. Such links must be left untouched - they should not be
# run through URL_REGEX again (that would produce broken/nested <a> tags).
EXISTING_LINK_REGEX = re.compile(
    r'<a\s[^>]*href\s*=\s*["\'][^"\']*["\'][^>]*>.*?</a>',
    re.IGNORECASE | re.DOTALL
)

# Regular expression catching literal formatting tags (b/i/u/strong/em/sub/sup)
# that may appear as plain text (e.g. pasted raw HTML), so we know when a text
# node needs to be re-parsed as real HTML rather than left as escaped text.
FORMATTING_TAG_REGEX = re.compile(r'</?(?:b|i|u|strong|em|sub|sup)\b[^>]*>', re.IGNORECASE)
 
class DocumentToHtmlConverter(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Link Converter: Word & PowerPoint -> HTML (Dark Theme)")
        self.resize(1000, 600)
 
        self.init_ui()
        self.apply_dark_theme()
 
    def init_ui(self):
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout(central_widget)
 
        # --- Button panel (Top) ---
        btn_layout = QHBoxLayout()
 
        self.btn_load = QPushButton("Load file (Word / PowerPoint)")
        self.btn_paste = QPushButton("Paste from clipboard")
        self.btn_copy = QPushButton("Copy resulting HTML")
 
        self.btn_load.clicked.connect(self.load_file)
        self.btn_paste.clicked.connect(self.paste_from_clipboard)
        self.btn_copy.clicked.connect(self.copy_result)
 
        btn_layout.addWidget(self.btn_load)
        btn_layout.addWidget(self.btn_paste)
        btn_layout.addWidget(self.btn_copy)
 
        main_layout.addLayout(btn_layout)
 
        # --- Text panel (Bottom - 2 columns) ---
        text_layout = QHBoxLayout()
 
        # Left side: Preview and Drag & Drop
        left_layout = QVBoxLayout()
        left_label = QLabel("Preview (You can drop text here or load a file):")
        self.preview_area = QTextEdit()
        self.preview_area.textChanged.connect(self.process_content)
 
        left_layout.addWidget(left_label)
        left_layout.addWidget(self.preview_area)
 
        # Right side: Resulting HTML code
        right_layout = QVBoxLayout()
        right_label = QLabel("Converted HTML code:")
        self.result_area = QPlainTextEdit()
        self.result_area.setReadOnly(True) 
 
        right_layout.addWidget(right_label)
        right_layout.addWidget(self.result_area)
 
        text_layout.addLayout(left_layout)
        text_layout.addLayout(right_layout)
 
        main_layout.addLayout(text_layout)
 
    def apply_dark_theme(self):
        dark_stylesheet = """
            QMainWindow, QWidget {
                background-color: #2b2b2b;
                color: #a9b7c6;
                font-family: 'Segoe UI', Arial, sans-serif;
            }
            QTextEdit, QPlainTextEdit {
                background-color: #1e1e1e;
                color: #a9b7c6;
                border: 1px solid #555555;
                border-radius: 6px;
                padding: 8px;
                font-size: 14px;
            }
            QPushButton {
                background-color: #365880;
                color: #ffffff;
                border: none;
                border-radius: 6px;
                padding: 10px 15px;
                font-weight: bold;
                font-size: 13px;
            }
            QPushButton:hover {
                background-color: #436a9b;
            }
            QPushButton:pressed {
                background-color: #27405e;
            }
            QLabel {
                font-size: 13px;
                font-weight: bold;
                margin-bottom: 5px;
                color: #cccccc;
            }
        """
        self.setStyleSheet(dark_stylesheet)
 
    def load_file(self):
        """Loads a .docx or .pptx file and processes it accordingly"""
        file_path, _ = QFileDialog.getOpenFileName(
            self, "Select file", "", "Documents (*.docx *.pptx)"
        )
        if not file_path:
            return
 
        try:
            if file_path.endswith('.docx'):
                with open(file_path, "rb") as docx_file:
                    # Map bold → <b>, italic → <i>, underline → <u>
                    # All three can coexist (mammoth wraps them independently).
                    style_map = (
                        "b => b\n"
                        "i => i\n"
                        "u => u\n"
                        "vertical-align('superscript') => sup\n"
                        "vertical-align('subscript') => sub"
                    )
                    result = mammoth.convert_to_html(docx_file, style_map=style_map)
                    self.preview_area.setHtml(result.value)
 
            elif file_path.endswith('.pptx'):
                html_content = self.parse_pptx(file_path)
                self.preview_area.setHtml(html_content)
 
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to load file:\n{str(e)}")
 
    def parse_pptx(self, file_path):
        """Extracts text, hyperlinks and formatting from PowerPoint text frames"""
        prs = Presentation(file_path)
        html_output = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    p_html = ""
                    for run in paragraph.runs:
                        text = html.escape(run.text)

                        # Wrap each formatting tag independently so they can stack:
                        # e.g. bold + underline → <b><u>text</u></b>
                        # Note: <u> is skipped for hyperlinks — browsers underline
                        # <a> elements by default, so it would be redundant.
                        font = run.font
                        is_link = bool(run.hyperlink and run.hyperlink.address)

                        # python-pptx has no direct superscript/subscript API,
                        # so read the raw <a:rPr baseline="..."/> attribute:
                        # positive baseline => superscript, negative => subscript.
                        rPr = run._r.find(qn('a:rPr'))
                        baseline = rPr.get('baseline') if rPr is not None else None
                        is_superscript = baseline is not None and int(baseline) > 0
                        is_subscript = baseline is not None and int(baseline) < 0

                        if is_superscript:
                            text = f"<sup>{text}</sup>"
                        if is_subscript:
                            text = f"<sub>{text}</sub>"
                        if font.underline and not is_link:
                            text = f"<u>{text}</u>"
                        if font.italic:
                            text = f"<i>{text}</i>"
                        if font.bold:
                            text = f"<b>{text}</b>"

                        if is_link:
                            clean_address = run.hyperlink.address.rstrip(',;.:')
                            address = html.escape(clean_address, quote=True)
                            p_html += f'<a href="{address}">{text}</a>'
                        else:
                            p_html += text


                    if p_html.strip():
                        html_output += f"<p>{p_html}</p>\n"

        return html_output
 
    def paste_from_clipboard(self):
        self.preview_area.paste()
 
    def copy_result(self):
        QApplication.clipboard().setText(self.result_area.toPlainText())
        #QMessageBox.information(self, "Success", "HTML code copied to clipboard!")
 
    def process_content(self):
        """Parses the preview code into clean HTML"""
        raw_html = self.preview_area.toHtml()
        soup = BeautifulSoup(raw_html, 'html.parser')

        body = soup.body if soup.body else soup

        # --- Step 1: Normalise Qt inline styles into semantic b/i/u tags ---
        # Qt encodes bold/italic/underline as inline CSS on <span> elements.
        # A single span can carry multiple styles at once, so we check all three
        # independently (no elif) and wrap the content in the appropriate tags.
        for span in body.find_all('span'):
            style = span.get('style', '')

            is_bold = ('font-weight:600' in style or 'font-weight: 600' in style or
                       'font-weight:700' in style or 'font-weight: 700' in style or
                       'font-weight:bold' in style or 'font-weight: bold' in style)
            is_italic = ('font-style:italic' in style or 'font-style: italic' in style)
            is_underline = ('text-decoration:underline' in style or
                            'text-decoration: underline' in style)
            is_superscript = ('vertical-align:super' in style or
                              'vertical-align: super' in style)
            is_subscript = ('vertical-align:sub' in style or
                            'vertical-align: sub' in style)

            if not (is_bold or is_italic or is_underline or is_superscript or is_subscript):
                continue

            # Collect the span's children, then wrap them in layers:
            # sup/sub → u → i → b
            # (innermost first so the outermost tag is the first one readers see)
            children = list(span.contents)

            if is_superscript:
                sup_tag = soup.new_tag('sup')
                for child in children:
                    sup_tag.append(child)
                children = [sup_tag]

            if is_subscript:
                sub_tag = soup.new_tag('sub')
                for child in children:
                    sub_tag.append(child)
                children = [sub_tag]

            if is_underline:
                u_tag = soup.new_tag('u')
                for child in children:
                    u_tag.append(child.__copy__() if hasattr(child, '__copy__') else child)
                children = [u_tag]

            if is_italic:
                i_tag = soup.new_tag('i')
                for child in children:
                    i_tag.append(child)
                children = [i_tag]

            if is_bold:
                b_tag = soup.new_tag('b')
                for child in children:
                    b_tag.append(child)
                children = [b_tag]

            span.replace_with(children[0])

        # --- Step 2: Convert raw http URLs in text nodes into <a> links ---
        # Also catches literal formatting/link tags pasted as plain text
        # (e.g. raw HTML source), re-parsing them into real elements so they
        # get picked up by the strong/em/u-in-a cleanup that follows.
        for text_node in body.find_all(string=True):
            if text_node.find_parent('a'):
                continue

            original_text = str(text_node)
            if ('http' not in original_text
                    and not EXISTING_LINK_REGEX.search(original_text)
                    and not FORMATTING_TAG_REGEX.search(original_text)):
                continue

            def clean_url_match(match):
                full_url = match.group(1)
                clean_url = full_url.rstrip(',;.:')
                return f'<a href="{clean_url}">{full_url}</a>'

            # The text node may already contain literal, ready-made <a href=...>
            # markup (e.g. pasted raw HTML). Cut those chunks out first so they
            # are left completely untouched, then run URL_REGEX only on the
            # remaining plain-text pieces to catch "naked" URLs.
            existing_links = EXISTING_LINK_REGEX.findall(original_text)
            plain_pieces = EXISTING_LINK_REGEX.split(original_text)

            rebuilt = ""
            for i, piece in enumerate(plain_pieces):
                rebuilt += URL_REGEX.sub(clean_url_match, piece)
                if i < len(existing_links):
                    rebuilt += existing_links[i]

            # Re-parse unconditionally: even if URL_REGEX made no change,
            # the text node may still contain literal formatting/link tags
            # (e.g. <b>...</b>) that need to become real HTML elements.
            new_soup = BeautifulSoup(rebuilt, 'html.parser')
            text_node.replace_with(new_soup)

        # --- Step 2b: Normalise <strong>/<em> into <b>/<i> ---
        # Runs after Step 2 on purpose, so it also catches <strong>/<em> that
        # came from literal HTML text re-parsed above (not just tags that
        # mammoth/Qt produced directly).
        for strong in body.find_all('strong'):
            strong.name = 'b'
        for em in body.find_all('em'):
            em.name = 'i'

        # --- Step 2c: Strip redundant <u> inside <a> ---
        # Browsers underline anchor elements by default, so <u> inside <a> is
        # purely noise. Runs after Step 2 so it also catches <u> that came
        # from literal HTML text re-parsed above. Unwrap every <u> that has
        # an <a> ancestor.
        for u_tag in body.find_all('u'):
            if u_tag.find_parent('a'):
                u_tag.unwrap()

        # --- Step 3: Build clean output — keep <a>, <b>, <i>, <u> ---
        KEEP_TAGS = {'a', 'b', 'i', 'u', 'sub', 'sup'}
        result_text = ""
        for block in body.find_all(['p', 'div', 'li', 'h1', 'h2', 'h3']):
            tags_to_unwrap = [tag for tag in block.find_all(True)
                              if tag.name not in KEEP_TAGS]
            for tag in tags_to_unwrap:
                tag.unwrap()

            block_html = "".join(str(c) for c in block.contents).strip()
            if block_html:
                result_text += block_html + "\n\n"

        self.result_area.setPlainText(result_text.strip())
 
if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = DocumentToHtmlConverter()
    window.show()
    sys.exit(app.exec())